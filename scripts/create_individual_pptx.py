#!/usr/bin/env python3
"""各社長ごとに個別のPPTX手順書を作成"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ===== Shared Colors =====
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
CARD_BG = RGBColor(0x2D, 0x2D, 0x44)
STEP_BG = RGBColor(0x23, 0x23, 0x3A)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_flat_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

def add_multiline_text(slide, left, top, width, height, lines, size=15, color=WHITE, line_spacing=1.3):
    """Add text with multiple lines, each as a paragraph"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(4)
    return txBox

def create_title_slide(prs, name, subtitle, accent_color, tagline):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Top accent bar
    add_flat_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.15), accent_color)

    # Name
    add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
        name, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Subtitle
    add_text(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.7),
        subtitle, size=22, color=accent_color, align=PP_ALIGN.CENTER)

    # Divider
    add_flat_rect(slide, Inches(5), Inches(4.1), Inches(3), Pt(2), accent_color)

    # Tagline
    add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8),
        tagline, size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Footer
    add_text(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
        "AI活用 実践手順書  |  clearAI", size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def create_overview_slide(prs, title, points, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
        title, size=30, color=WHITE, bold=True)
    add_flat_rect(slide, Inches(0.6), Inches(1.1), Inches(3), Pt(3), accent_color)

    for i, (label, desc) in enumerate(points):
        y = Inches(1.5) + i * Inches(1.2)
        card = add_rect(slide, Inches(0.6), y, Inches(12), Inches(1.0), CARD_BG)
        # Number
        circle = add_circle(slide, Inches(0.9), y + Inches(0.15), Inches(0.6), accent_color)
        tf = circle.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text(slide, Inches(1.7), y + Inches(0.05), Inches(10.5), Inches(0.4),
            label, size=18, color=WHITE, bold=True)
        add_text(slide, Inches(1.7), y + Inches(0.45), Inches(10.5), Inches(0.5),
            desc, size=14, color=LIGHT_GRAY)

def create_step_slide(prs, step_num, step_title, tool_name, steps, tips, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Header
    header = add_flat_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.3), accent_color)
    add_text(slide, Inches(0.6), Inches(0.15), Inches(2), Inches(0.5),
        f"STEP {step_num}", size=16, color=WHITE, bold=True)
    add_text(slide, Inches(0.6), Inches(0.55), Inches(10), Inches(0.6),
        step_title, size=28, color=WHITE, bold=True)

    # Tool badge
    add_text(slide, Inches(9), Inches(0.15), Inches(4), Inches(0.5),
        f"使用ツール: {tool_name}", size=14, color=WHITE, align=PP_ALIGN.RIGHT)

    # Steps (left side)
    add_text(slide, Inches(0.6), Inches(1.5), Inches(4), Inches(0.4),
        "手順", size=16, color=accent_color, bold=True)

    for i, step in enumerate(steps):
        y = Inches(2.0) + i * Inches(0.85)
        card = add_rect(slide, Inches(0.6), y, Inches(7.2), Inches(0.7), STEP_BG)
        # Step number circle
        c = add_circle(slide, Inches(0.8), y + Inches(0.1), Inches(0.45), RGBColor(0x3D, 0x3D, 0x5C))
        tf = c.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(16)
        p.font.color.rgb = accent_color
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        add_text(slide, Inches(1.5), y + Inches(0.1), Inches(6.0), Inches(0.55),
            step, size=14, color=WHITE)

    # Tips (right side)
    add_text(slide, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.4),
        "ポイント・コツ", size=16, color=RGBColor(0xFF, 0xD7, 0x00), bold=True)
    tips_card = add_rect(slide, Inches(8.3), Inches(2.0), Inches(4.5), Inches(4.8), CARD_BG)
    add_multiline_text(slide, Inches(8.5), Inches(2.2), Inches(4.1), Inches(4.4),
        tips, size=13, color=LIGHT_GRAY)


def create_closing_slide(prs, name, quote, accent_color):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.7),
        f"{name}流 AI活用のまとめ", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_flat_rect(slide, Inches(5), Inches(2.3), Inches(3), Pt(2), accent_color)

    # Quote box
    quote_card = add_rect(slide, Inches(2), Inches(3.0), Inches(9), Inches(2.0), CARD_BG)
    add_text(slide, Inches(2.3), Inches(3.1), Inches(1), Inches(0.8),
        "\"", size=60, color=accent_color, bold=True)
    add_text(slide, Inches(3.0), Inches(3.5), Inches(7.5), Inches(1.2),
        quote, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
        "まずは1つのツールから始めてみましょう！\nご相談は clearAI まで",
        size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def save_pptx(prs, filename):
    output_dir = os.path.dirname(os.path.dirname(__file__))
    path = os.path.join(output_dir, filename)
    prs.save(path)
    print(f"Saved: {path}")


# ============================================================
# 1. 南場智子 - Perplexity × NotebookLM × Circleback
# ============================================================
def create_nanba():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0x00, 0x96, 0xFF)

    create_title_slide(prs,
        "南場 智子（DeNA会長）のAI活用術",
        "Perplexity × NotebookLM × Circleback で情報戦を制する",
        ACCENT,
        "「経営者がAIに興奮しているかがポイント」")

    create_overview_slide(prs, "南場流 AI活用の全体像", [
        ("Perplexity AI で人物リサーチ", "初めて会う人の「必読記事」をAIに聞き、即座に重要情報を収集"),
        ("NotebookLM で瞬時にインプット", "記事URL・YouTube動画を全部放り込み、チャットで深掘り質問"),
        ("Circleback で会議を自動記録", "ミーティングの会話をAIが自動で議事録化。振り返りも瞬時に"),
        ("Deep Research で投資判断", "ChatGPTのDeep Research機能で投資先の徹底調査を短時間で完了"),
    ], ACCENT)

    create_step_slide(prs, 1,
        "Perplexity AIで「必読記事」を見つける",
        "Perplexity AI（perplexity.ai）",
        [
            "Perplexity AI（perplexity.ai）を開く",
            "「〇〇さんについての必読記事は何ですか？」と入力",
            "表示された記事のURLをすべてコピーする",
            "X（Twitter）の発信もチェック → URLを控える",
            "YouTubeで最近の登壇動画があれば、そのURLも控える",
        ],
        [
            "Perplexityは出典付きで回答するので",
            "信頼性の高い記事が見つかりやすい",
            "",
            "「最近1年以内の記事に絞って」と",
            "追加指示すると、より新鮮な情報が得られる",
            "",
            "Pro版（月$20）ならDeep Researchも使える",
        ], ACCENT)

    create_step_slide(prs, 2,
        "NotebookLMに全情報を投入 → チャットで深掘り",
        "Google NotebookLM（notebooklm.google.com）",
        [
            "NotebookLM（notebooklm.google.com）を開く",
            "「新しいノートブック」を作成する",
            "STEP1で集めたURLを「ソースを追加」で全部入れる",
            "YouTube動画のURLもそのまま入れる（自動文字起こし）",
            "チャット欄で「この方の最近の関心事は？」と質問する",
        ],
        [
            "南場氏は「行きのタクシーの中で」",
            "チャットで質問すると語っている",
            "",
            "「トランプ政権についてどう考えて",
            "いるか」など具体的な質問が効果的",
            "",
            "NotebookLMは無料で使える！",
            "Googleアカウントがあればすぐ開始",
        ], ACCENT)

    create_step_slide(prs, 3,
        "Circlebackで会議の議事録を自動作成",
        "Circleback（circleback.ai）",
        [
            "Circleback（circleback.ai）でアカウント作成",
            "Google Calendar / Outlook と連携設定する",
            "会議が始まると自動でAIが録音・文字起こし開始",
            "会議終了後、自動で要約・アクションアイテムが生成される",
            "後日「あの会議で〇〇について何と言っていた？」と検索可能",
        ],
        [
            "Zoom / Google Meet / Teams に対応",
            "",
            "日本語の会議にも対応している",
            "",
            "「誰が何を言ったか」も記録されるので",
            "責任の所在が明確になる",
            "",
            "無料プランでも月数回は使える",
        ], ACCENT)

    create_step_slide(prs, 4,
        "Deep Researchで投資・事業判断の調査",
        "ChatGPT Deep Research",
        [
            "ChatGPT（Pro/Plus）を開く",
            "モデル選択で「Deep Research」を選ぶ",
            "「〇〇市場の最新動向と主要プレイヤーを調査して」と入力",
            "AIが数分〜数十分かけて100以上のソースを自動で調査",
            "出典付きの詳細なレポートが生成される",
        ],
        [
            "南場氏は投資判断にこの機能を活用",
            "",
            "従来なら調査チームが数日かかる作業が",
            "数十分で完了する",
            "",
            "ChatGPT Plus（月$20）以上で利用可能",
            "",
            "レポートはPDFでエクスポートもできる",
        ], ACCENT)

    create_closing_slide(prs, "南場智子",
        "経営者がAIに興奮しているかがポイント", ACCENT)

    save_pptx(prs, "手順書_南場智子_AI活用術.pptx")


# ============================================================
# 2. 堀江貴文 - ChatGPTでコンテンツ作成・意思決定
# ============================================================
def create_horie():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0xFF, 0x6B, 0x35)

    create_title_slide(prs,
        "堀江 貴文（ホリエモン）のAI活用術",
        "ChatGPTを「もう一人の自分」として使い倒す",
        ACCENT,
        "「今やらないヤツはバカ！」")

    create_overview_slide(prs, "堀江流 AI活用の全体像", [
        ("ChatGPTで書籍を丸ごと作成", "自分の音声データをAIに渡し、99%をAIが書いたビジネス書を出版"),
        ("スピーチ・挨拶文を瞬時に生成", "イベント登壇時のスピーチ原稿をChatGPTで自動生成"),
        ("マイGPT（カスタムGPT）で専門AIを作る", "自社の就業規則を読み込ませた「総務部長GPT」など、業務特化AIを構築"),
        ("日常のアイデア出し・壁打ち相手", "新規事業のアイデアや判断をChatGPTに相談し高速でPDCAを回す"),
    ], ACCENT)

    create_step_slide(prs, 1,
        "ChatGPTでスピーチ・文章を一発生成",
        "ChatGPT（chat.openai.com）",
        [
            "ChatGPT（chat.openai.com）を開く",
            "「以下の条件でスピーチ原稿を作って」と指示",
            "条件を具体的に書く（場面・対象者・時間・トーン）",
            "生成された原稿を読み、修正点を追加指示",
            "「もっとカジュアルに」「数字を入れて」等で磨き上げ",
        ],
        [
            "堀江氏は本の帯コメントもAIで生成",
            "",
            "ポイントは「条件を具体的に書くこと」",
            "・誰に向けた話か",
            "・何分の話か",
            "・フォーマルかカジュアルか",
            "",
            "最初から完璧を求めず、",
            "何度もやり取りして磨くのがコツ",
        ], ACCENT)

    create_step_slide(prs, 2,
        "マイGPT（カスタムGPT）で社内専門AIを作る",
        "ChatGPT Plus/Team（GPTs機能）",
        [
            "ChatGPT左メニューの「GPTsを探す」→「作成」をクリック",
            "「あなたは何をするGPTですか？」に役割を入力",
            "「ナレッジ」に自社の資料（PDF/Word等）をアップロード",
            "テスト質問で動作確認 → 修正を繰り返す",
            "完成したら社内メンバーにリンクを共有",
        ],
        [
            "堀江氏は就業規則を読み込ませて",
            "「総務部長GPT」を作ることを推奨",
            "",
            "他にも作れる例：",
            "・営業マニュアルGPT",
            "・商品知識GPT",
            "・クレーム対応GPT",
            "",
            "ChatGPT Plus（月$20）で利用可能",
        ], ACCENT)

    create_step_slide(prs, 3,
        "AIで本を1冊書く（コンテンツ量産）",
        "ChatGPT + 音声入力",
        [
            "自分の考えをスマホの音声メモで録音する",
            "Whisper等の文字起こしツールでテキスト化",
            "ChatGPTに「以下の文字起こしを章立てして書籍にして」と依頼",
            "各章ごとに「もっと具体例を入れて」等で品質を上げる",
            "最終チェックは自分の目で行い、出版・社内共有",
        ],
        [
            "堀江氏は実際に99%をAIが書いた",
            "ビジネス書『夢を叶える力』を出版",
            "",
            "社長の考えを社内に共有するための",
            "「社長の考え方ブック」制作にも応用可能",
            "",
            "Whisperは無料で使える文字起こしAI",
            "（openai.com/research/whisper）",
        ], ACCENT)

    create_step_slide(prs, 4,
        "新規事業のアイデア壁打ち",
        "ChatGPT（GPT-4o / o1）",
        [
            "ChatGPTを開き「新規事業のアイデアを一緒に考えて」と入力",
            "自社の強み・リソース・課題を箇条書きで伝える",
            "出てきたアイデアに「なぜそれが上手くいくと思う？」と深掘り",
            "「リスクと対策を3つずつ挙げて」で現実チェック",
            "良いアイデアは「事業計画書の叩き台を作って」で具体化",
        ],
        [
            "AIは「いい質問をする力」で結果が変わる",
            "",
            "漠然と「何かいいアイデアない？」ではなく",
            "具体的な制約条件を与えるのがコツ",
            "",
            "例：「予算500万円以内で」",
            "「3ヶ月で立ち上げられる」",
            "「既存の顧客基盤を活かせる」",
        ], ACCENT)

    create_closing_slide(prs, "堀江貴文",
        "今やらないヤツはバカ！", ACCENT)

    save_pptx(prs, "手順書_堀江貴文_AI活用術.pptx")


# ============================================================
# 3. Sam Altman - ChatGPT Pulse・メール処理・生活のAI化
# ============================================================
def create_altman():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0x10, 0xA3, 0x7F)

    create_title_slide(prs,
        "Sam Altman（OpenAI CEO）のAI活用術",
        "「退屈なタスク」こそAIに任せる ― 日常を変える使い方",
        ACCENT,
        "「AIは退屈な仕事にこそ使え」")

    create_overview_slide(prs, "Altman流 AI活用の全体像", [
        ("メール処理を丸ごとAIに任せる", "大量のメールの処理・要約・返信下書きをChatGPTで高速化"),
        ("文書の要約で情報処理を加速", "長文レポートや契約書の要約を一瞬で作成"),
        ("ChatGPT Pulseで毎朝パーソナル更新", "夜間にAIが自動分析、朝に自分向けのブリーフィングを受け取る"),
        ("子育て・日常生活の相談相手", "育児の疑問からレシピ提案まで、生活のあらゆる場面でAI活用"),
    ], ACCENT)

    create_step_slide(prs, 1,
        "大量メールの処理をAIで高速化",
        "ChatGPT（chat.openai.com）",
        [
            "処理したいメール本文をコピーする",
            "ChatGPTに「以下のメールを要約して、返信が必要なものをリストアップして」",
            "要約を見て優先順位を判断する",
            "「このメールへの返信を、丁寧に断るトーンで書いて」と依頼",
            "生成された返信を確認 → 微修正してメールに貼り付け",
        ],
        [
            "Altmanは「退屈な仕事にこそAIを使う」",
            "と明言している",
            "",
            "メール処理は最も効果が出やすい活用法",
            "",
            "Gmail連携のChatGPTプラグインや",
            "Copilot for Outlookもおすすめ",
            "",
            "まずは1日10通から始めてみよう",
        ], ACCENT)

    create_step_slide(prs, 2,
        "長文文書の要約を一瞬で作成",
        "ChatGPT（GPT-4o）",
        [
            "要約したい文書（PDF/Word等）をChatGPTにアップロード",
            "「この文書を3行で要約して」と指示",
            "「重要なポイントを箇条書きで5つ抽出して」と追加指示",
            "「この内容について質問：〇〇はどうなっている？」で深掘り",
            "「経営会議向けの1ページサマリーを作って」で資料化",
        ],
        [
            "契約書、調査レポート、競合資料など",
            "読む時間がない文書に最適",
            "",
            "PDFは直接アップロードできる",
            "（ChatGPT Plus以上）",
            "",
            "「専門用語を避けてわかりやすく」",
            "と追加指示すると、社長向けに最適化される",
        ], ACCENT)

    create_step_slide(prs, 3,
        "ChatGPT Pulseで毎朝のブリーフィング",
        "ChatGPT Pulse（最新機能）",
        [
            "ChatGPTアプリを開き、プロフィール設定から「メモリ」をONにする",
            "日頃から関心事をChatGPTに話しかけておく（業界ニュース等）",
            "Pulse機能をONにする（設定 → Pulse）",
            "毎朝、Pulseが「あなた向けの更新情報」を自動生成",
            "関心事の進展、前日の会話のフォローアップ等が届く",
        ],
        [
            "Altmanが「最もお気に入りの機能」と語る",
            "",
            "夜間にAIが自動でネットを巡回し、",
            "自分の関心事に関連する情報を整理",
            "",
            "「秘書が毎朝ブリーフィング」を",
            "AIが自動で実現するイメージ",
            "",
            "ChatGPT Plus/Proで利用可能",
        ], ACCENT)

    create_step_slide(prs, 4,
        "日常生活の「ちょっとした疑問」にAIを使う",
        "ChatGPTアプリ（スマートフォン）",
        [
            "スマホにChatGPTアプリをインストール",
            "音声モードをONにする（マイクアイコンをタップ）",
            "話しかけるだけで回答が返ってくる",
            "「今夜のレシピ提案して」「子供にこう聞かれたら？」等",
            "メモリ機能で好みを覚えてもらい、どんどんパーソナルに",
        ],
        [
            "Altmanは育児の質問にもAIを使う",
            "",
            "音声モードなら運転中や料理中も",
            "ハンズフリーで使える",
            "",
            "「メモリ」機能で家族の食の好みや",
            "アレルギー情報を覚えさせると",
            "提案の精度が格段に上がる",
        ], ACCENT)

    create_closing_slide(prs, "Sam Altman",
        "AIは退屈な仕事にこそ使え。それが一番効果が出る", ACCENT)

    save_pptx(prs, "手順書_SamAltman_AI活用術.pptx")


# ============================================================
# 4. Tobi Lutke (Shopify) - Cursor × Claude × Copilot
# ============================================================
def create_lutke():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0x96, 0xBF, 0x48)

    create_title_slide(prs,
        "Tobi Lutke（Shopify CEO）のAI活用術",
        "「AIで100倍の成果」を出す組織の作り方",
        ACCENT,
        "「AI活用は全社員の基本的な期待事項」")

    create_overview_slide(prs, "Lutke流 AI活用の全体像", [
        ("AIをプロトタイプの相棒にする", "新機能のアイデアをAIで即座にプロトタイプ化、意思決定を加速"),
        ("「採用前にAIを試す」ルール", "新規採用の前に「AIで代替できないか？」を証明することを義務化"),
        ("AI活用度を人事評価に組み込む", "全社員のパフォーマンスレビューにAIスキルを評価項目として追加"),
        ("Claude / Cursor / Copilotの使い分け", "用途に応じて最適なAIツールを選択するフレームワーク"),
    ], ACCENT)

    create_step_slide(prs, 1,
        "アイデアをAIで即プロトタイプ化",
        "Claude（claude.ai）/ ChatGPT",
        [
            "新しいアイデアが浮かんだらすぐにAIを開く",
            "「〇〇のプロトタイプを作りたい。要件は以下：」と入力",
            "要件を箇条書きで列挙（ユーザー・機能・制約）",
            "AIが出したプロトタイプ案に対して「ここをもっと簡単に」等で修正",
            "プロトタイプ案が固まったらチームに共有して議論",
        ],
        [
            "Lutkeは「GSD（Get Shit Done）",
            "プロトタイプフェーズ」でのAI活用を重視",
            "",
            "完璧なものを作る前に、",
            "AIで素早く叩き台を作ることが重要",
            "",
            "Claudeの「Artifacts」機能なら",
            "HTMLプロトタイプも作れる",
        ], ACCENT)

    create_step_slide(prs, 2,
        "「この仕事はAIでできない？」チェック",
        "社内ルールとして導入",
        [
            "新規採用リクエストが来たら、まず「AIで代替できないか？」と問う",
            "具体的に「どのタスクをAIに任せたか？」を報告させる",
            "AIで80%解決できるなら、残り20%を既存メンバーに振り分け",
            "それでも人が必要な場合のみ採用を承認",
            "この判断プロセスを四半期ごとに見直す",
        ],
        [
            "Shopifyでは全チームにこのルールを適用",
            "",
            "「AIでできることをやらせずに",
            "人を雇うのはムダ」という考え方",
            "",
            "経営者から始めることが重要！",
            "自分が先にAIを使い倒す姿を見せる",
        ], ACCENT)

    create_step_slide(prs, 3,
        "AIツールを用途別に使い分ける",
        "Claude / Cursor / GitHub Copilot",
        [
            "長文の分析・戦略立案 → Claude（200Kトークンの長文対応）",
            "コード開発・自動化 → Cursor（AIコードエディタ）",
            "日常のコーディング補助 → GitHub Copilot",
            "情報検索・調査 → Perplexity AI",
            "各ツールの得意分野を把握し、タスクに応じて切り替える",
        ],
        [
            "Shopifyは全社員にこれらのツールを提供",
            "",
            "まずは1つから始めて、",
            "慣れたら用途別に使い分けを覚える",
            "",
            "Claude → 思考のパートナー",
            "Cursor → 開発のパートナー",
            "Copilot → 日常の効率化",
        ], ACCENT)

    create_step_slide(prs, 4,
        "人事評価にAI活用度を組み込む",
        "人事制度の設計",
        [
            "パフォーマンスレビューの評価項目に「AI活用度」を追加",
            "「AIをどのように業務に活用したか」を自己申告させる",
            "チームリーダーがAI活用の好事例を共有する場を設ける",
            "AI活用で成果を上げた社員を表彰する仕組みを作る",
            "四半期ごとにAI活用の成果をチーム単位で振り返る",
        ],
        [
            "「使え」と言うだけでは変わらない",
            "評価に組み込むことで行動が変わる",
            "",
            "Lutkeは「100倍の成果を出した社員」",
            "の事例を全社に共有している",
            "",
            "まずは小さな成功事例を積み上げて",
            "社内の「AIは使える」という空気を作る",
        ], ACCENT)

    create_closing_slide(prs, "Tobi Lutke",
        "AI活用はもはや選択肢ではない。全員の基本的な期待事項だ", ACCENT)

    save_pptx(prs, "手順書_TobiLutke_AI活用術.pptx")


# ============================================================
# 5. 落合陽一 - Claude Code × vibe-local × Notion
# ============================================================
def create_ochiai():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0x9B, 0x59, 0xB6)

    create_title_slide(prs,
        "落合 陽一（ピクシーダストCEO）のAI活用術",
        "Claude Codeに「寝てる間に作らせる」最先端の使い方",
        ACCENT,
        "「2026年にはほとんどの知的作業がAIに置き換わる」")

    create_overview_slide(prs, "落合流 AI活用の全体像", [
        ("Claude Codeで「寝ている間に」開発", "夜にAIに指示を出して寝ると、朝にはツールが完成している"),
        ("Notionで全情報をWiki化", "申請書の書き方からWi-Fi設定まで、あらゆる情報をNotionに集約"),
        ("「審美眼」でAI出力を選別する", "膨大なAI出力から「良いもの」を選ぶ力こそが人間の価値"),
        ("vibe-local でオフラインAIコーディング", "ネット不要・サブスク不要の完全ローカルAI開発環境"),
    ], ACCENT)

    create_step_slide(prs, 1,
        "Claude Codeに「寝ている間に」タスクを実行させる",
        "Claude Code（ターミナル）",
        [
            "ターミナル（コマンドライン）でClaude Codeを起動",
            "作りたいものを自然言語で詳しく指示する",
            "「完了したらファイルを保存して」と追加指示",
            "AIが自動でコードを書き始める → そのまま就寝",
            "朝起きて結果を確認 → 修正が必要なら追加指示",
        ],
        [
            "落合氏は実際にClaude Codeに",
            "「Claude Codeを作って」と指示して",
            "寝たら、朝にはツールが完成していた",
            "",
            "経営者の場合、Webサイトや",
            "社内ツールの作成にも応用可能",
            "",
            "プログラミング知識がなくても",
            "日本語の指示で動く",
        ], ACCENT)

    create_step_slide(prs, 2,
        "Notionで「自分だけのWiki」を構築",
        "Notion（notion.so）",
        [
            "Notion（notion.so）でワークスペースを作成",
            "「社内マニュアル」「議事録」「プロジェクト」等のページを作る",
            "あらゆる情報をNotionに集約する（URLリンクも可）",
            "Notion AI機能で「この文書を要約して」と指示できる",
            "チームメンバーと共有し、全員が同じ情報にアクセス",
        ],
        [
            "落合氏は1人で50人の学生を見ており",
            "「Notionがないと仕事が回らない」と語る",
            "",
            "社長が持つ暗黙知を",
            "Notionに言語化して蓄積すると",
            "組織の知識資産になる",
            "",
            "無料プランでも十分に使える",
        ], ACCENT)

    create_step_slide(prs, 3,
        "AI出力の「審美眼」を鍛える",
        "Claude / ChatGPT / 各種AIツール",
        [
            "同じ質問を複数のAI（ChatGPT, Claude, Gemini）に投げる",
            "それぞれの回答を比較し、最も優れたものを選ぶ",
            "「なぜこの回答が良いと思ったか」を自分で言語化する",
            "AIの出力を「そのまま使わず」必ず自分の視点で編集する",
            "この選別作業を繰り返すことで「審美眼」が磨かれる",
        ],
        [
            "落合氏：AIを上手に使える人と",
            "使えない人の差は「審美眼」",
            "",
            "つまり「選ぶ力」が重要",
            "",
            "AIが100個のアイデアを出したら",
            "その中から最良の1個を選べるかが勝負",
            "",
            "経営判断に直結するスキル",
        ], ACCENT)

    create_step_slide(prs, 4,
        "vibe-localでネット不要のAI環境を構築",
        "vibe-local（落合氏開発ツール）",
        [
            "GitHub から vibe-local をダウンロード",
            "ローカルLLM（Ollama等）をインストール",
            "vibe-localを起動すると、ネット不要でAIコーディングが可能",
            "機密情報を含むプロジェクトでも安心して使える",
            "社内ネットワーク外でもAIの力を活用できる",
        ],
        [
            "完全無料・ネットワーク不要",
            "",
            "機密情報がクラウドに送信されない",
            "ため、セキュリティ面で安心",
            "",
            "技術者向けだが、IT部門に",
            "依頼すれば社内展開も可能",
            "",
            "外部にデータを出せない業界に最適",
        ], ACCENT)

    create_closing_slide(prs, "落合陽一",
        "2026年にはほとんどの知的作業がAIに置き換わる", ACCENT)

    save_pptx(prs, "手順書_落合陽一_AI活用術.pptx")


# ============================================================
# 6. Mustafa Suleyman - Copilot で個人の生産性を極限まで高める
# ============================================================
def create_suleyman():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    ACCENT = RGBColor(0x00, 0x78, 0xD4)

    create_title_slide(prs,
        "Mustafa Suleyman（Microsoft AI CEO）のAI活用術",
        "Copilotを「AIコンパニオン」として日常に溶け込ませる",
        ACCENT,
        "「AIに任せられる知的作業は、もうほぼ全てだ」")

    create_overview_slide(prs, "Suleyman流 AI活用の全体像", [
        ("Copilotを個人データベースとして使う", "映画リスト、読書メモ、個人の好みをAIに記憶させて活用"),
        ("Copilot for Microsoft 365で業務効率化", "Word/Excel/Teams/Outlookの全てにAIを統合し業務を高速化"),
        ("「AIコンパニオン」として長期記憶を活用", "過去の会話を記憶し、自分を深く理解するAIパートナーを育てる"),
        ("あらゆる知的作業をAIに委ねる", "「どんな知識作業でもAIに聞いてみる」を習慣化する"),
    ], ACCENT)

    create_step_slide(prs, 1,
        "Copilotを「個人データベース」にする",
        "Microsoft Copilot（copilot.microsoft.com）",
        [
            "Copilot（copilot.microsoft.com）を開く",
            "「メモリ」機能をONにする（設定から）",
            "好きな映画・本・レストラン等をCopilotに伝える",
            "「次に見るべき映画を推薦して」と聞くとパーソナルな提案",
            "仕事の好みや判断基準も伝えておくと業務でも活用",
        ],
        [
            "Suleyman自身が映画のパーソナル",
            "データベースとしてCopilotを活用",
            "",
            "日付、感想、評価をAIに記録させ",
            "「あの時見た映画なんだっけ？」にも",
            "答えてもらえる",
            "",
            "経営の意思決定基準を記録させると",
            "「先月の判断と矛盾してない？」も聞ける",
        ], ACCENT)

    create_step_slide(prs, 2,
        "Copilot for Microsoft 365 で業務を高速化",
        "Microsoft 365 + Copilot",
        [
            "Microsoft 365管理者がCopilotライセンスを有効化",
            "Wordで「このメモから企画書を作って」→ 一瞬で文書生成",
            "Excelで「このデータの傾向を分析して」→ グラフ付きレポート",
            "Teamsで「この会議の要約とアクションアイテムを」→ 自動生成",
            "Outlookで「このメールスレッドを要約して」→ 重要ポイント抽出",
        ],
        [
            "既にWord/Excel/Teamsを使っているなら",
            "最も導入ハードルが低いAI",
            "",
            "特にTeamsの会議要約は",
            "「参加できなかった会議の内容を",
            "5分で把握できる」と好評",
            "",
            "月額3,750円/ユーザー（税別）",
        ], ACCENT)

    create_step_slide(prs, 3,
        "AIに「自分のことを覚えてもらう」",
        "Copilot / ChatGPT（メモリ機能）",
        [
            "AIとの会話で、自分の好み・判断基準・役割を伝える",
            "「私は〇〇の社長で、△△業界に詳しい」と自己紹介する",
            "仕事の進め方の好みも伝える（データ重視 / 直感重視 等）",
            "過去の会話を踏まえた提案が返ってくるようになる",
            "定期的に「私について何を覚えている？」で確認・修正",
        ],
        [
            "Suleymanのビジョン：",
            "「全ての人にAIコンパニオンを」",
            "",
            "AIに自分を理解させるほど",
            "提案の質が劇的に向上する",
            "",
            "「過去3ヶ月の議論を踏まえて、",
            "次のステップを提案して」等が可能に",
        ], ACCENT)

    create_step_slide(prs, 4,
        "「とりあえずAIに聞いてみる」を習慣化",
        "Copilot / ChatGPT / Claude（全般）",
        [
            "迷ったらまずAIに聞く →「〇〇についてどう思う？」",
            "知らないことがあったらAIに聞く →「〇〇を教えて」",
            "文章を書く前にAIに叩き台を作らせる",
            "スケジュール調整もAIに相談 →「来週の優先順位を整理して」",
            "1日に最低10回はAIに話しかけることを目標にする",
        ],
        [
            "Suleymanは「プロフェッショナルの",
            "タスクの大部分は12〜18ヶ月以内に",
            "AIで自動化される」と予測",
            "",
            "今から「AIに聞く習慣」を身につけて",
            "おくことが重要",
            "",
            "最初は質問の質が低くてもOK",
            "使い続けることで上手くなる",
        ], ACCENT)

    create_closing_slide(prs, "Mustafa Suleyman",
        "あらゆる知的作業をAIに聞いてみよう。それが未来の働き方だ", ACCENT)

    save_pptx(prs, "手順書_MustafaSuleyman_AI活用術.pptx")


# ============================================================
# メイン実行
# ============================================================
if __name__ == "__main__":
    create_nanba()
    create_horie()
    create_altman()
    create_lutke()
    create_ochiai()
    create_suleyman()
    print("\n全6ファイルの作成が完了しました！")
