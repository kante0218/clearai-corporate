#!/usr/bin/env python3
"""AI活用で有名な社長6名のPowerPoint資料を作成"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== Color Palette =====
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xFF)
ACCENT_PURPLE = RGBColor(0x7C, 0x3A, 0xED)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
CARD_BG = RGBColor(0x2D, 0x2D, 0x44)
HIGHLIGHT = RGBColor(0x00, 0xD4, 0xAA)
ORANGE = RGBColor(0xFF, 0x8C, 0x00)

def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return txBox

# ===== CEO Data =====
ceos = [
    {
        "name": "南場 智子",
        "company": "DeNA（ディー・エヌ・エー）",
        "title": "代表取締役会長",
        "country": "日本",
        "flag": "JP",
        "color": ACCENT_BLUE,
        "slogan": "「AIにオールイン」宣言",
        "overview": "2025年に「DeNAはAIにオールイン」を宣言。\n経営者自ら生成AIのヘビーユーザーとして\n具体的な活用法を公開し、大きな話題に。",
        "tools": [
            "Perplexity AI → 会議前の情報収集・人物リサーチ",
            "NotebookLM → 大量資料の要約・分析",
            "ChatGPT o1 → 戦略立案・プロトタイピング",
            "Circleback → 会議の自動議事録作成",
        ],
        "results": [
            "QA業務の工数を50%削減",
            "Pococha配信審査コスト60%削減",
            "現業3,000名→半分でも発展する体制へ",
            "バーティカルAIエージェント事業を新規立ち上げ",
        ],
        "quote": "経営者がAIに興奮しているかがポイント",
    },
    {
        "name": "孫 正義",
        "company": "ソフトバンクグループ",
        "title": "代表取締役会長兼社長",
        "country": "日本",
        "flag": "JP",
        "color": ACCENT_PURPLE,
        "slogan": "ASI時代のNo.1プラットフォーマーへ",
        "overview": "AI時代の覇者を目指し、OpenAIに4.8兆円投資。\nグループ全体で10億AIエージェント構想を発表。\n「Stargate Project」で世界最大級のAIインフラ構築。",
        "tools": [
            "Crystal Intelligence → 自己増殖型AIエージェント群",
            "社員1人あたり1,000個のAIエージェント常時稼働",
            "ネットワーク監視～人事・財務・営業まで全部門AI化",
            "生成AIコンテスト → 5万人から累計26万件のアイデア",
        ],
        "results": [
            "OpenAIに4.8兆円投資（歴史上最大規模）",
            "Stargate Projectで米国AIインフラ構築",
            "グループ全体で10億AIエージェント計画",
            "ASI（人工超知能）プラットフォーム構想",
        ],
        "quote": "数千兆のAIエージェントが常時稼働する世界が来る",
    },
    {
        "name": "三木谷 浩史",
        "company": "楽天グループ",
        "title": "代表取締役会長兼社長",
        "country": "日本",
        "flag": "JP",
        "color": RGBColor(0xBF, 0x00, 0x00),
        "slogan": "「No AI, No Future」宣言",
        "overview": "楽天市場を世界最先端のAI活用ECプラットフォームへ。\n「AIの民主化」を掲げ、全従業員のAI活用を推進。\n2030年に楽天市場の流通総額10兆円を目指す。",
        "tools": [
            "Rakuten AI → 独自開発の生成AIを楽天市場に導入",
            "OpenAI連携 → 戦略的パートナーシップ",
            "トリプル20 → マーケ・オペレーション・顧客効率各20%向上",
            "全従業員3万人がAI活用、毎日8,000人が利用",
        ],
        "results": [
            "マーケティング効率20%向上を達成",
            "オペレーション効率20%向上を目標",
            "2030年流通総額10兆円へ向けた成長加速",
            "「AIの民主化」で全社員のスキル底上げ",
        ],
        "quote": "AIを使いこなせない企業・人の未来は厳しい",
    },
    {
        "name": "Satya Nadella",
        "company": "Microsoft",
        "title": "CEO（最高経営責任者）",
        "country": "アメリカ",
        "flag": "US",
        "color": RGBColor(0x00, 0x78, 0xD4),
        "slogan": "AIを「認知増幅装置」として全製品に統合",
        "overview": "Copilotブランドの下、Office・Azure・GitHub等\nあらゆる製品にAIを統合。OpenAIへの巨額投資を主導し\nエンタープライズAIの世界標準を構築。",
        "tools": [
            "Microsoft Copilot → Office全製品にAIアシスタント統合",
            "Azure OpenAI Service → 企業向けAIプラットフォーム",
            "GitHub Copilot → 開発者の生産性を劇的に向上",
            "Copilot+ PC → AIネイティブなハードウェア展開",
        ],
        "results": [
            "AI事業が年間売上130億ドル超を達成",
            "Azure成長率がAI需要で加速",
            "GitHub Copilotが開発時間55%短縮",
            "2026年を「AIの転換年」と位置づけ",
        ],
        "quote": "AIは人間の能力を置き換えるのではなく増幅する",
    },
    {
        "name": "Marc Benioff",
        "company": "Salesforce",
        "title": "会長兼CEO",
        "country": "アメリカ",
        "flag": "US",
        "color": RGBColor(0x00, 0xA1, 0xE0),
        "slogan": "社内業務の50%をAIが担当",
        "overview": "Agentforceプラットフォームで「デジタル従業員」を推進。\n社内業務の30〜50%をAIが処理する体制を実現。\nダボス会議で「全員が人間だけの組織を率いる最後の世代」と発言。",
        "tools": [
            "Agentforce → AIエージェントプラットフォーム",
            "顧客対応の85%をAIが自動解決",
            "営業リード選別を40%高速化",
            "エンジニアリング・マーケティング全般をAI支援",
        ],
        "results": [
            "社内業務の30〜50%をAIが処理",
            "CS部門4,000名分の業務をAIが代替",
            "2030年に年間売上600億ドルを目標",
            "51%の採用が社内リデプロイメント",
        ],
        "quote": "我々は全員が人間だけの組織を率いる最後のCEO世代だ",
    },
    {
        "name": "Tobi Lutke",
        "company": "Shopify",
        "title": "CEO（最高経営責任者）",
        "country": "カナダ",
        "flag": "CA",
        "color": RGBColor(0x96, 0xBF, 0x48),
        "slogan": "「AI活用は全社員の基本的な期待事項」",
        "overview": "2025年に全社員向けメモでAI活用を義務化。\n新規採用の前に「AIで代替できない」証明を要求。\n人事評価にもAI活用度を組み込むAI-first組織を構築。",
        "tools": [
            "全社員にAI活用を「基本的な期待事項」として義務化",
            "新規採用前にAI代替不可の証明を義務づけ",
            "人事評価にAI活用度を組み込み",
            "Shopify Magic → EC事業者向けAIツール群",
        ],
        "results": [
            "組織全体のAI-first文化を確立",
            "不要な新規採用を抑制しコスト最適化",
            "EC事業者にAIツールを提供し差別化",
            "AI活用を人事評価の重要指標に設定",
        ],
        "quote": "AIを使えないことは、もはや選択肢ではない",
    },
]

# ===== Slide 1: Title =====
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Title
add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
    "AI活用で世界をリードする経営者たち",
    font_size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.0), Inches(11), Inches(0.8),
    "〜 6人のトップリーダーに学ぶ、AI経営の最前線 〜",
    font_size=24, color=ACCENT_BLUE, bold=False, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
    "2025-2026年 最新動向まとめ",
    font_size=18, color=LIGHT_GRAY, bold=False, align=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(2.8), Inches(4), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

# Footer
add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
    "Presented by clearAI",
    font_size=14, color=LIGHT_GRAY, bold=False, align=PP_ALIGN.CENTER)


# ===== Slide 2: Overview (6 CEOs at a glance) =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "AI活用で注目される6人の経営者",
    font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Decorative line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.1), Inches(5), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

card_w = Inches(3.8)
card_h = Inches(2.6)
gap_x = Inches(0.35)
start_x = Inches(0.65)
row_y = [Inches(1.5), Inches(4.3)]

for i, ceo in enumerate(ceos):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = row_y[row]

    card = add_shape_bg(slide, x, y, card_w, card_h, CARD_BG)

    # Accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Pt(5), card_h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ceo["color"]
    bar.line.fill.background()

    add_text_box(slide, x + Inches(0.3), y + Inches(0.15), card_w - Inches(0.4), Inches(0.5),
        ceo["name"], font_size=20, color=WHITE, bold=True)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), card_w - Inches(0.4), Inches(0.4),
        f"{ceo['company']}  |  {ceo['country']}", font_size=13, color=LIGHT_GRAY)

    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), card_w - Inches(0.4), Inches(0.4),
        ceo["slogan"], font_size=14, color=HIGHLIGHT, bold=True)

    # Short quote
    add_text_box(slide, x + Inches(0.3), y + Inches(1.5), card_w - Inches(0.4), Inches(0.9),
        f"「{ceo['quote']}」", font_size=11, color=LIGHT_GRAY)


# ===== Individual CEO Slides =====
for ceo in ceos:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Header bar
    header_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(1.5))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = ceo["color"]
    header_bar.line.fill.background()

    add_text_box(slide, Inches(0.7), Inches(0.2), Inches(8), Inches(0.7),
        ceo["name"], font_size=36, color=WHITE, bold=True)
    add_text_box(slide, Inches(0.7), Inches(0.85), Inches(8), Inches(0.5),
        f"{ceo['company']}  {ceo['title']}  |  {ceo['country']}",
        font_size=16, color=RGBColor(0xFF, 0xFF, 0xFF))

    # Slogan badge
    add_text_box(slide, Inches(8.5), Inches(0.35), Inches(4.5), Inches(0.8),
        ceo["slogan"], font_size=18, color=WHITE, bold=True, align=PP_ALIGN.RIGHT)

    # Overview section
    add_text_box(slide, Inches(0.5), Inches(1.7), Inches(3), Inches(0.4),
        "OVERVIEW", font_size=14, color=ACCENT_BLUE, bold=True)
    overview_card = add_shape_bg(slide, Inches(0.5), Inches(2.1), Inches(5.8), Inches(2.0), CARD_BG)
    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(5.3), Inches(1.8),
        ceo["overview"], font_size=15, color=WHITE)

    # Tools / AI活用 section
    add_text_box(slide, Inches(6.8), Inches(1.7), Inches(3), Inches(0.4),
        "AI活用ツール・施策", font_size=14, color=ACCENT_BLUE, bold=True)
    tools_card = add_shape_bg(slide, Inches(6.8), Inches(2.1), Inches(6.0), Inches(2.0), CARD_BG)
    tools_text = "\n".join([f"  {t}" for t in ceo["tools"]])
    add_text_box(slide, Inches(7.0), Inches(2.2), Inches(5.6), Inches(1.8),
        tools_text, font_size=14, color=WHITE)

    # Results section
    add_text_box(slide, Inches(0.5), Inches(4.3), Inches(3), Inches(0.4),
        "主な成果・インパクト", font_size=14, color=HIGHLIGHT, bold=True)
    results_card = add_shape_bg(slide, Inches(0.5), Inches(4.7), Inches(5.8), Inches(2.3), CARD_BG)
    results_text = "\n".join([f"  {r}" for r in ceo["results"]])
    add_text_box(slide, Inches(0.8), Inches(4.8), Inches(5.3), Inches(2.1),
        results_text, font_size=15, color=WHITE)

    # Quote section
    add_text_box(slide, Inches(6.8), Inches(4.3), Inches(3), Inches(0.4),
        "注目の発言", font_size=14, color=ORANGE, bold=True)
    quote_card = add_shape_bg(slide, Inches(6.8), Inches(4.7), Inches(6.0), Inches(2.3), CARD_BG)

    # Quote mark
    add_text_box(slide, Inches(7.0), Inches(4.9), Inches(1), Inches(0.8),
        "\"", font_size=60, color=ceo["color"], bold=True)
    add_text_box(slide, Inches(7.5), Inches(5.3), Inches(5.0), Inches(1.5),
        ceo["quote"], font_size=20, color=WHITE, bold=True)


# ===== Summary Slide =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
    "AI経営の共通ポイント",
    font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.1), Inches(5), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

points = [
    ("01", "経営者自身がAIのヘビーユーザーである", "南場氏のように経営者自らAIを使い倒し、具体的な活用法を社員に示すことで組織全体のAI活用が加速する。"),
    ("02", "AI前提で組織・業務を再設計する", "単にツールを導入するだけでなく、Benioff氏やLutke氏のように人事制度・採用基準・業務フローをAI前提で再構築する。"),
    ("03", "全社員をAI人材化する仕組みを作る", "三木谷氏の「AIの民主化」、孫氏の「生成AIコンテスト」のように、全社員がAIを使える環境と動機を整備する。"),
    ("04", "AIエージェントが次のフロンティア", "6名全員がAIエージェント（自律的に業務を遂行するAI）を次の主戦場と位置づけている。"),
]

colors_pts = [ACCENT_BLUE, ACCENT_PURPLE, HIGHLIGHT, ORANGE]

for i, (num, title_text, desc) in enumerate(points):
    y = Inches(1.5) + i * Inches(1.45)
    card = add_shape_bg(slide, Inches(0.8), y, Inches(11.5), Inches(1.25), CARD_BG)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.15), Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = colors_pts[i]
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

    add_text_box(slide, Inches(2.0), y + Inches(0.1), Inches(10), Inches(0.45),
        title_text, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.55), Inches(10), Inches(0.6),
        desc, font_size=14, color=LIGHT_GRAY)


# ===== Last Slide =====
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
    "AI経営は「未来」ではなく「今」",
    font_size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.0),
    "トップリーダーたちは既に動いています。\n御社のAI経営戦略、一緒に考えませんか？",
    font_size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.2), Inches(3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT_BLUE
line.line.fill.background()

add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    "clearAI  |  AI導入コンサルティング",
    font_size=18, color=ACCENT_BLUE, bold=True, align=PP_ALIGN.CENTER)

# Save
output_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "AI活用リーダー6名.pptx")
prs.save(output_path)
print(f"PowerPoint saved: {output_path}")
